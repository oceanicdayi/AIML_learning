"""
extract_pptx.py - 提取 PPTX 簡報的所有文字與圖片資訊

使用方式:
    python scripts/extract_pptx.py <pptx_path>
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_slide_content(pptx_path: str) -> list[dict]:
    prs = Presentation(pptx_path)
    slides_data = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_info = {
            "slide_number": slide_num,
            "layout": slide.slide_layout.name if slide.slide_layout else "Unknown",
            "texts": [],
            "tables": [],
            "images": [],
            "notes": "",
        }

        # Extract notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            slide_info["notes"] = slide.notes_slide.notes_text_frame.text.strip()

        for shape in slide.shapes:
            # Text frames
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        level = para.level
                        slide_info["texts"].append({"text": text, "level": level})

            # Tables
            if shape.has_table:
                table = shape.table
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                slide_info["tables"].append(table_data)

            # Images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = shape.image
                    slide_info["images"].append({
                        "content_type": img.content_type,
                        "size": f"{shape.width}x{shape.height}",
                        "name": shape.name,
                    })
                except ValueError:
                    slide_info["images"].append({
                        "content_type": "linked",
                        "size": f"{shape.width}x{shape.height}",
                        "name": shape.name,
                    })

            # Group shapes (recursive)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for s in shape.shapes:
                    if hasattr(s, "text_frame") and s.has_text_frame:
                        for para in s.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_info["texts"].append({"text": text, "level": para.level})

        slides_data.append(slide_info)

    return slides_data


def format_output(slides: list[dict]) -> str:
    lines = []
    for slide in slides:
        lines.append(f"\n{'='*60}")
        lines.append(f"## Slide {slide['slide_number']} (Layout: {slide['layout']})")
        lines.append(f"{'='*60}")

        for item in slide["texts"]:
            prefix = "  " * item["level"] + "- " if item["level"] > 0 else ""
            lines.append(f"{prefix}{item['text']}")

        for table in slide["tables"]:
            lines.append("\n[TABLE]")
            for row in table:
                lines.append(" | ".join(row))
            lines.append("[/TABLE]")

        if slide["images"]:
            lines.append(f"\n[IMAGES: {len(slide['images'])} image(s)]")
            for img in slide["images"]:
                lines.append(f"  - {img['name']} ({img['content_type']})")

        if slide["notes"]:
            lines.append(f"\n[NOTES]: {slide['notes']}")

    return "\n".join(lines)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_pptx.py <pptx_path>")
        sys.exit(1)

    pptx_path = sys.argv[1]
    slides = extract_slide_content(pptx_path)
    print(format_output(slides))
    print(f"\n\nTotal slides: {len(slides)}")
